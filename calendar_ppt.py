from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from calendar import monthcalendar, setfirstweekday, SUNDAY
from holidays import Taiwan

# 設定每週的第一天為星期日
setfirstweekday(SUNDAY)

# 獲取台灣法定假日
official_holidays = Taiwan(years=2025, language="zh_TW")

# 手動添加其他假日或生日
additional_holidays = {
    (1, 9): "捏生日",
    (2, 12): "元宵節",
    (2, 14): "情人節",
    (4, 2): "玲生日",
    (5, 1): "勞動節",
    (5, 11): "母親節",
    (5, 22): "88生日",
    (10, 6): "中秋節",
    (10, 28): "蓮生日",
    (12, 25): "聖誕節"
}

# 合併假日
taiwan_holidays = {}
for date, name in official_holidays.items():
    month, day = date.month, date.day
    if month not in taiwan_holidays:
        taiwan_holidays[month] = {}
    taiwan_holidays[month][day] = name

for (month, day), name in additional_holidays.items():
    if month not in taiwan_holidays:
        taiwan_holidays[month] = {}
    taiwan_holidays[month][day] = name

# 簡單的月份名稱
month_names = [
    "一月", "二月", "三月", "四月", "五月", "六月",
    "七月", "八月", "九月", "十月", "十一月", "十二月"
]

# 建立PPT
prs = Presentation()
slide_width = Cm(38)
slide_height = Cm(52)
prs.slide_width = slide_width
prs.slide_height = slide_height
year = 2025

# 生成每個月的日曆頁面
for month in range(1, 13):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版面
    shapes = slide.shapes

    # 插入上方圖片的區域（此處僅預留）
    pic_placeholder = shapes.add_shape(
        1, Cm(0.5), Cm(0.5), slide_width - Cm(1), Cm(20)  # 圖片區域佔 4:6 比例的 4
    )

    # 插入月曆表格的標題，並使其居中
    title = shapes.add_textbox(Cm(14), Cm(20.5), slide_width - Cm(2), Cm(2))  # 放置在圖片區和表格區之間
    title_frame = title.text_frame
    title_frame.text = f"{year}年 {month_names[month-1]}"
    title_frame.paragraphs[0].font.size = Pt(55)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = 1  # 使標題居中

    # 設定日曆表格
    table = shapes.add_table(rows=7, cols=7, left=Cm(1), top=Cm(23), width=slide_width - Cm(2), height=Cm(28)).table  # 表格區佔 6

    # 表格內容填入
    days = ["日", "一", "二", "三", "四", "五", "六"]
    for col, day in enumerate(days):
        table.cell(0, col).text = day
        table.cell(0, col).text_frame.paragraphs[0].font.bold = True
        table.cell(0, col).text_frame.paragraphs[0].font.size = Pt(60)  # 星期字體放大

    cal = monthcalendar(year, month)
    for row_idx, week in enumerate(cal):
        for col_idx, day in enumerate(week):
            cell = table.cell(row_idx + 1, col_idx)
            if day != 0:  # 填入日期
                cell.text = str(day)
                # 設定日期字體大小
                cell.text_frame.paragraphs[0].font.size = Pt(60)

                # 判斷是否為週末，並將日期字體顏色設為紅色
                if col_idx == 0 or col_idx == 6:  # 星期日和星期六
                    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
                
                # 判斷是否是假日，並將日期字體顏色設為紅色
                if month in taiwan_holidays and day in taiwan_holidays[month]:
                    holiday_name = taiwan_holidays[month][day]
                    
                    # 判斷是否為額外假日（橘色），否則為法定假日（紅色）
                    if (month, day) in additional_holidays:
                        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 165, 0)  # 橘色
                    else:
                        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)  # 紅色
                    
                    # 添加假日名稱為新的段落
                    p = cell.text_frame.add_paragraph()
                    p.text = holiday_name
                    p.font.size = Pt(30)  # 設定字體大小
                    p.font.color.rgb = cell.text_frame.paragraphs[0].font.color.rgb  # 保持顏色一致

# 儲存PPT檔案
output_path = "台灣日曆.pptx"
prs.save(output_path)
print(f"日曆已成功生成，保存於：{output_path}")
